"""
Inspect a .pptx template to list all slide layouts and their placeholders.
Usage: python scripts/inspect_template.py [path/to/template.pptx]
"""
import sys
import os
from pptx import Presentation
from pptx.util import Emu


def inspect(path: str) -> None:
    prs = Presentation(path)

    print(f"File      : {path}")
    print(f"Slide size: {Emu(prs.slide_width).inches:.2f}\" x {Emu(prs.slide_height).inches:.2f}\"")
    print(f"Layouts   : {len(prs.slide_layouts)}\n")

    for idx, layout in enumerate(prs.slide_layouts):
        print(f"[{idx}] \"{layout.name}\"")
        if not layout.placeholders:
            print("     (no placeholders)")
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            print(
                f"     idx={ph.placeholder_format.idx}"
                f"  type={ph_type.name if ph_type else 'UNKNOWN':<20}"
                f"  name=\"{ph.name}\""
            )
        print()


if __name__ == "__main__":
    default_path = os.path.join(
        os.path.dirname(__file__), "..", "static", "template.pptx"
    )
    template_path = sys.argv[1] if len(sys.argv) > 1 else default_path
    inspect(os.path.normpath(template_path))
