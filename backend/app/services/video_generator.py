"""
Video generation pipeline:
  source content → slide JSON (Groq) → PPTX → PNG images (LibreOffice) →
  parallel [image decoration (Gemini) + TTS (Gemini)] → MP4 (moviepy)
"""
from __future__ import annotations

import json
import logging
import os
import re
import subprocess
import wave
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from typing import TYPE_CHECKING, Union

from google import genai
from google.genai import types
from langchain_groq import ChatGroq
from langchain_postgres import PGVector
from pptx import Presentation
from pptx.util import Inches

from app.core.config import settings
from app.services.embedding_service import get_embeddings

if TYPE_CHECKING:
    from app.models.video import Video

logger = logging.getLogger(__name__)

# ------------------------------------------------------------------
# Slide dataclasses — one per layout
# ------------------------------------------------------------------

@dataclass
class TitleSlide:
    title: str
    subtitle: str
    speaker_notes: str = ""


@dataclass
class TitleAndContentSlide:
    title: str
    bullets: list[str] = field(default_factory=list)
    speaker_notes: str = ""


@dataclass
class SectionHeaderSlide:
    title: str
    text: str = ""
    speaker_notes: str = ""


@dataclass
class TwoContentSlide:
    title: str
    bullets_left: list[str] = field(default_factory=list)
    bullets_right: list[str] = field(default_factory=list)
    speaker_notes: str = ""


@dataclass
class ComparisonSlide:
    title: str
    header_left: str = ""
    bullets_left: list[str] = field(default_factory=list)
    header_right: str = ""
    bullets_right: list[str] = field(default_factory=list)
    speaker_notes: str = ""


@dataclass
class TitleOnlySlide:
    title: str
    speaker_notes: str = ""


@dataclass
class ContentWithCaptionSlide:
    title: str
    bullets: list[str] = field(default_factory=list)
    caption: str = ""
    speaker_notes: str = ""


@dataclass
class PictureWithCaptionSlide:
    title: str
    image_prompt: str = ""
    caption: str = ""
    speaker_notes: str = ""


SlideData = Union[
    TitleSlide,
    TitleAndContentSlide,
    SectionHeaderSlide,
    TwoContentSlide,
    ComparisonSlide,
    TitleOnlySlide,
    ContentWithCaptionSlide,
    PictureWithCaptionSlide,
]

# Maps JSON "type" value → (dataclass, layout name)
_SLIDE_TYPE_MAP: dict[str, tuple[type, str]] = {
    "title_slide":           (TitleSlide,              "Title Slide"),
    "title_and_content":     (TitleAndContentSlide,    "Title and Content"),
    "section_header":        (SectionHeaderSlide,      "Section Header"),
    "two_content":           (TwoContentSlide,         "Two Content"),
    "comparison":            (ComparisonSlide,         "Comparison"),
    "title_only":            (TitleOnlySlide,          "Title Only"),
    "content_with_caption":  (ContentWithCaptionSlide, "Content with Caption"),
    "picture_with_caption":  (PictureWithCaptionSlide, "Picture with Caption"),
}

_LLM_SCHEMA = """
Available slide types and their JSON fields:

- "title_slide"          : {"type", "title", "subtitle", "speaker_notes"}
- "title_and_content"    : {"type", "title", "bullets": [...], "speaker_notes"}
- "section_header"       : {"type", "title", "text", "speaker_notes"}
- "two_content"          : {"type", "title", "bullets_left": [...], "bullets_right": [...], "speaker_notes"}
- "comparison"           : {"type", "title", "header_left", "bullets_left": [...], "header_right", "bullets_right": [...], "speaker_notes"}
- "title_only"           : {"type", "title", "speaker_notes"}
- "content_with_caption" : {"type", "title", "bullets": [...], "caption", "speaker_notes"}
- "picture_with_caption" : {"type", "title", "image_prompt", "caption", "speaker_notes"}
"""


# ------------------------------------------------------------------
# Main generator class
# ------------------------------------------------------------------

class VideoGenerator:
    def __init__(self, video: "Video", notebook_id: str, output_dir: str) -> None:
        self.video = video
        self.notebook_id = notebook_id
        self.output_dir = output_dir
        self.tmp_dir = os.path.join(output_dir, "tmp")
        os.makedirs(self.tmp_dir, exist_ok=True)

        self._genai_client = genai.Client(api_key=settings.GOOGLE_API_KEY)
        self._llm = ChatGroq(
            api_key=settings.GROQ_API_KEY,
            model=settings.GROQ_MODEL,
            temperature=0.3,
        )

    def generate(self) -> str:
        """Run the full pipeline. Returns path to output .mp4 file."""
        logger.info("Step 1: Fetching source content")
        content = self._get_source_content()

        logger.info("Step 2: Generating slide content via LLM")
        slides = self._generate_slide_content(content)

        logger.info("Step 3: Building PPTX")
        pptx_path = self._build_pptx(slides)

        logger.info("Step 4: Exporting PPTX to images")
        image_paths = self._export_to_images(pptx_path)

        # Align slides and images (LibreOffice should produce 1 PNG per slide)
        count = min(len(slides), len(image_paths))
        slides = slides[:count]
        image_paths = image_paths[:count]

        logger.info("Step 5: Parallel image decoration + TTS")
        decorated_images, audio_paths = self._parallel_process(slides, image_paths)

        logger.info("Step 6: Rendering video")
        mp4_path = self._render_video(decorated_images, audio_paths)

        return mp4_path

    # ------------------------------------------------------------------
    # Step 1: Fetch source content from PGVector
    # ------------------------------------------------------------------

    def _get_source_content(self) -> str:
        from sqlalchemy import create_engine
        from sqlalchemy.orm import Session

        from app.models.source import Source

        engine = create_engine(settings.SYNC_DATABASE_URL)
        embeddings = get_embeddings()

        with Session(engine) as db:
            sources = (
                db.query(Source)
                .filter(
                    Source.notebook_id == self.notebook_id,
                    Source.status == "ready",
                )
                .all()
            )

        all_docs = []
        for source in sources:
            store = PGVector(
                embeddings=embeddings,
                collection_name=f"source_{source.id}",
                connection=settings.SYNC_DATABASE_URL,
            )
            docs = store.similarity_search("key concepts and main points", k=5)
            all_docs.extend(docs)

        if not all_docs:
            return "No content available."

        return "\n\n---\n\n".join(doc.page_content for doc in all_docs)

    # ------------------------------------------------------------------
    # Step 2: Generate slide structure via Groq LLM
    # ------------------------------------------------------------------

    def _generate_slide_content(self, content: str) -> list[SlideData]:
        prompt = (
            "You are a presentation designer. Based on the following source content, "
            "create a structured slide deck with 6-10 slides.\n\n"
            "Rules:\n"
            "- First slide MUST be type 'title_slide'\n"
            "- Last slide MUST be type 'section_header' or 'title_only' (as closing)\n"
            "- Use a variety of slide types to keep the presentation engaging\n"
            "- Use 'comparison' or 'two_content' when contrasting two ideas\n"
            "- Use 'picture_with_caption' for concepts that benefit from a visual; "
            "  write a detailed, descriptive image_prompt for Gemini image generation\n"
            "- speaker_notes should be 2-3 natural sentences suitable for narration\n\n"
            f"{_LLM_SCHEMA}\n"
            'Return ONLY a valid JSON object: {"slides": [...]}\n\n'
            f"SOURCE CONTENT:\n{content[:8000]}"
        )

        response = self._llm.invoke(prompt)
        raw = response.content if hasattr(response, "content") else str(response)

        json_match = re.search(r"\{.*\}", raw, re.DOTALL)
        if not json_match:
            raise ValueError(f"LLM did not return valid JSON: {raw[:200]}")

        data = json.loads(json_match.group())
        slides: list[SlideData] = []
        for s in data.get("slides", []):
            slide = _parse_slide(s)
            if slide is not None:
                slides.append(slide)
        return slides

    # ------------------------------------------------------------------
    # Step 3: Build PPTX from slide data
    # ------------------------------------------------------------------

    def _build_pptx(self, slides: list[SlideData]) -> str:
        template_path = settings.VIDEO_TEMPLATE_PATH
        if os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            logger.warning("Template not found at %s, using blank presentation", template_path)
            prs = Presentation()

        # Index layouts by name for reliable lookup
        layouts: dict[str, object] = {l.name: l for l in prs.slide_layouts}

        # Remove any placeholder slides that came with the template
        for _ in range(len(prs.slides)):
            rId = prs.slides._sldIdLst[0].get("r:id")
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        for slide_data in slides:
            _, layout_name = _SLIDE_TYPE_MAP[_slide_type_key(slide_data)]
            layout = layouts.get(layout_name)
            if layout is None:
                logger.warning("Layout '%s' not found in template, skipping slide", layout_name)
                continue

            slide = prs.slides.add_slide(layout)
            self._fill_slide(slide, slide_data)

            if slide_data.speaker_notes:
                slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

        pptx_path = os.path.join(self.tmp_dir, "presentation.pptx")
        prs.save(pptx_path)
        return pptx_path

    def _fill_slide(self, slide: object, data: SlideData) -> None:
        ph = {p.placeholder_format.idx: p for p in slide.placeholders}

        if isinstance(data, TitleSlide):
            _set_text(ph, 0, data.title)
            _set_text(ph, 1, data.subtitle)

        elif isinstance(data, TitleAndContentSlide):
            _set_text(ph, 0, data.title)
            _set_bullets(ph, 1, data.bullets)

        elif isinstance(data, SectionHeaderSlide):
            _set_text(ph, 0, data.title)
            _set_text(ph, 1, data.text)

        elif isinstance(data, TwoContentSlide):
            _set_text(ph, 0, data.title)
            _set_bullets(ph, 1, data.bullets_left)
            _set_bullets(ph, 2, data.bullets_right)

        elif isinstance(data, ComparisonSlide):
            _set_text(ph, 0, data.title)
            _set_text(ph, 1, data.header_left)
            _set_bullets(ph, 2, data.bullets_left)
            _set_text(ph, 3, data.header_right)
            _set_bullets(ph, 4, data.bullets_right)

        elif isinstance(data, TitleOnlySlide):
            _set_text(ph, 0, data.title)

        elif isinstance(data, ContentWithCaptionSlide):
            _set_text(ph, 0, data.title)
            _set_bullets(ph, 1, data.bullets)
            _set_text(ph, 2, data.caption)

        elif isinstance(data, PictureWithCaptionSlide):
            _set_text(ph, 0, data.title)
            _set_text(ph, 2, data.caption)
            # Picture placeholder (idx=1) filled after image generation in decoration step

    # ------------------------------------------------------------------
    # Step 4: Export PPTX to PNG images using LibreOffice
    # ------------------------------------------------------------------

    def _export_to_images(self, pptx_path: str) -> list[str]:
        images_dir = os.path.join(self.tmp_dir, "slides")
        os.makedirs(images_dir, exist_ok=True)

        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", images_dir,
                pptx_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        png_files = sorted(
            f for f in os.listdir(images_dir) if f.lower().endswith(".png")
        )
        return [os.path.join(images_dir, f) for f in png_files]

    # ------------------------------------------------------------------
    # Step 5: Parallel image decoration + TTS
    # ------------------------------------------------------------------

    def _parallel_process(
        self, slides: list[SlideData], image_paths: list[str]
    ) -> tuple[list[str | None], list[str | None]]:
        decorated: list[str | None] = [None] * len(slides)
        audios: list[str | None] = [None] * len(slides)

        with ThreadPoolExecutor(max_workers=4) as executor:
            decor_futures = {
                executor.submit(self._decorate_image, img_path, slides[i], i): i
                for i, img_path in enumerate(image_paths)
            }
            tts_futures = {
                executor.submit(self._generate_tts, slides[i].speaker_notes, i): i
                for i in range(len(slides))
            }

            for future in as_completed(decor_futures):
                idx = decor_futures[future]
                try:
                    decorated[idx] = future.result()
                except Exception as exc:
                    logger.warning("Image decoration failed for slide %d: %s", idx, exc)
                    decorated[idx] = image_paths[idx]

            for future in as_completed(tts_futures):
                idx = tts_futures[future]
                try:
                    audios[idx] = future.result()
                except Exception as exc:
                    logger.warning("TTS failed for slide %d: %s", idx, exc)
                    audios[idx] = None

        return decorated, audios

    def _decorate_image(self, img_path: str, slide: SlideData, idx: int) -> str:
        # PictureWithCaption: generate image from prompt, composite onto slide PNG
        if isinstance(slide, PictureWithCaptionSlide) and slide.image_prompt:
            return self._generate_picture_slide(img_path, slide, idx)

        # All other slides: add handwritten-style annotations
        with open(img_path, "rb") as f:
            image_bytes = f.read()

        prompt = (
            "Add human-like handwritten annotations, arrows, and simple illustrations "
            "to enhance this presentation slide. Keep the original text readable. "
            "Add only light, helpful annotations that a teacher might draw on a whiteboard. "
            f"The slide title is: {slide.title}"
        )

        response = self._genai_client.models.generate_content(
            model="gemini-2.0-flash-preview-image-generation",
            contents=[
                types.Part.from_bytes(data=image_bytes, mime_type="image/png"),
                prompt,
            ],
            config=types.GenerateContentConfig(
                response_modalities=["IMAGE"],
            ),
        )

        decorated_path = os.path.join(self.tmp_dir, f"decorated_{idx:03d}.png")
        for part in response.candidates[0].content.parts:
            if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                with open(decorated_path, "wb") as f:
                    f.write(part.inline_data.data)
                return decorated_path

        return img_path

    def _generate_picture_slide(
        self, slide_png_path: str, slide: PictureWithCaptionSlide, idx: int
    ) -> str:
        """Generate an image from slide.image_prompt and composite it onto the slide PNG."""
        from PIL import Image as PILImage

        # Generate image via Gemini
        response = self._genai_client.models.generate_content(
            model="gemini-2.0-flash-preview-image-generation",
            contents=slide.image_prompt,
            config=types.GenerateContentConfig(
                response_modalities=["IMAGE"],
            ),
        )

        generated_path = os.path.join(self.tmp_dir, f"generated_{idx:03d}.png")
        for part in response.candidates[0].content.parts:
            if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                with open(generated_path, "wb") as f:
                    f.write(part.inline_data.data)
                break
        else:
            return slide_png_path  # fallback if generation failed

        # Composite: paste generated image into the picture area of the slide PNG
        slide_img = PILImage.open(slide_png_path).convert("RGBA")
        gen_img = PILImage.open(generated_path).convert("RGBA")

        # Estimate picture area: roughly center-left, 55% width, 65% height
        sw, sh = slide_img.size
        pic_w = int(sw * 0.55)
        pic_h = int(sh * 0.65)
        pic_x = int(sw * 0.03)
        pic_y = int(sh * 0.22)

        gen_img = gen_img.resize((pic_w, pic_h), PILImage.LANCZOS)
        slide_img.paste(gen_img, (pic_x, pic_y))

        output_path = os.path.join(self.tmp_dir, f"decorated_{idx:03d}.png")
        slide_img.convert("RGB").save(output_path)
        return output_path

    def _generate_tts(self, speaker_notes: str, idx: int) -> str | None:
        if not speaker_notes.strip():
            return None

        audio_path = os.path.join(self.tmp_dir, f"audio_{idx:03d}.wav")

        response = self._genai_client.models.generate_content(
            model="gemini-2.5-flash-preview-tts",
            contents=speaker_notes,
            config=types.GenerateContentConfig(
                response_modalities=["AUDIO"],
                speech_config=types.SpeechConfig(
                    voice_config=types.VoiceConfig(
                        prebuilt_voice_config=types.PrebuiltVoiceConfig(
                            voice_name=self.video.voice_name or "Kore",
                        )
                    )
                ),
            ),
        )

        pcm_data = b""
        for part in response.candidates[0].content.parts:
            if part.inline_data and part.inline_data.mime_type.startswith("audio/"):
                pcm_data = part.inline_data.data
                break

        if not pcm_data:
            return None

        _save_pcm_as_wav(pcm_data, audio_path)
        return audio_path

    # ------------------------------------------------------------------
    # Step 6: Render final MP4 using moviepy
    # ------------------------------------------------------------------

    def _render_video(
        self, image_paths: list[str | None], audio_paths: list[str | None]
    ) -> str:
        from moviepy import AudioFileClip, ImageClip, concatenate_videoclips

        DEFAULT_SLIDE_DURATION = 5.0

        clips = []
        for img_path, audio_path in zip(image_paths, audio_paths):
            if img_path is None:
                continue
            if audio_path and os.path.exists(audio_path):
                audio_clip = AudioFileClip(audio_path)
                video_clip = (
                    ImageClip(img_path)
                    .with_duration(audio_clip.duration)
                    .with_audio(audio_clip)
                )
            else:
                video_clip = ImageClip(img_path).with_duration(DEFAULT_SLIDE_DURATION)
            clips.append(video_clip)

        if not clips:
            raise ValueError("No clips to render — pipeline produced no output")

        final = concatenate_videoclips(clips, method="compose")
        output_path = os.path.join(self.output_dir, "output.mp4")
        final.write_videofile(
            output_path,
            fps=24,
            codec="libx264",
            audio_codec="aac",
            logger=None,
        )
        return output_path


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------

def _slide_type_key(slide: SlideData) -> str:
    return next(k for k, (cls, _) in _SLIDE_TYPE_MAP.items() if isinstance(slide, cls))


def _parse_slide(s: dict) -> SlideData | None:
    type_key = s.get("type", "")
    entry = _SLIDE_TYPE_MAP.get(type_key)
    if entry is None:
        logger.warning("Unknown slide type '%s', skipping", type_key)
        return None

    cls, _ = entry
    match cls:
        case _ if cls is TitleSlide:
            return TitleSlide(
                title=s.get("title", ""),
                subtitle=s.get("subtitle", ""),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is TitleAndContentSlide:
            return TitleAndContentSlide(
                title=s.get("title", ""),
                bullets=s.get("bullets", []),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is SectionHeaderSlide:
            return SectionHeaderSlide(
                title=s.get("title", ""),
                text=s.get("text", ""),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is TwoContentSlide:
            return TwoContentSlide(
                title=s.get("title", ""),
                bullets_left=s.get("bullets_left", []),
                bullets_right=s.get("bullets_right", []),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is ComparisonSlide:
            return ComparisonSlide(
                title=s.get("title", ""),
                header_left=s.get("header_left", ""),
                bullets_left=s.get("bullets_left", []),
                header_right=s.get("header_right", ""),
                bullets_right=s.get("bullets_right", []),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is TitleOnlySlide:
            return TitleOnlySlide(
                title=s.get("title", ""),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is ContentWithCaptionSlide:
            return ContentWithCaptionSlide(
                title=s.get("title", ""),
                bullets=s.get("bullets", []),
                caption=s.get("caption", ""),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _ if cls is PictureWithCaptionSlide:
            return PictureWithCaptionSlide(
                title=s.get("title", ""),
                image_prompt=s.get("image_prompt", ""),
                caption=s.get("caption", ""),
                speaker_notes=s.get("speaker_notes", ""),
            )
        case _:
            return None


def _set_text(ph: dict, idx: int, text: str) -> None:
    if idx in ph and text:
        ph[idx].text = text


def _set_bullets(ph: dict, idx: int, bullets: list[str]) -> None:
    if idx not in ph or not bullets:
        return
    tf = ph[idx].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
        else:
            tf.add_paragraph().text = bullet


def _save_pcm_as_wav(pcm_data: bytes, path: str, sample_rate: int = 24000) -> None:
    with wave.open(path, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(sample_rate)
        wav_file.writeframes(pcm_data)
