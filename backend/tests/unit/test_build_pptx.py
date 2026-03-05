"""
Unit tests for VideoGenerator._build_pptx.

These tests call _build_pptx directly — no DB, no LLM, no network.
The real template (static/template.pptx) is used so layout names and
placeholder indices are verified against the actual file.
"""
import os
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation

from app.services.video_generator import (
    ComparisonSlide,
    ContentWithCaptionSlide,
    PictureWithCaptionSlide,
    SectionHeaderSlide,
    TitleAndContentSlide,
    TitleOnlySlide,
    TitleSlide,
    TwoContentSlide,
    VideoGenerator,
)

TEMPLATE_PATH = Path(__file__).parent.parent / "static" / "template.pptx"


@pytest.fixture
def generator(tmp_path: Path) -> VideoGenerator:
    video = MagicMock()
    video.voice_name = "Kore"
    with (
        patch("app.services.video_generator.genai.Client"),
        patch("app.services.video_generator.ChatGroq"),
    ):
        gen = VideoGenerator(
            video=video,
            notebook_id="00000000-0000-0000-0000-000000000001",
            output_dir=str(tmp_path),
        )
    return gen


def build(generator: VideoGenerator, slides: list) -> Presentation:
    """Call _build_pptx with an overridden template path and return the loaded result."""
    import app.core.config as cfg_module
    original = cfg_module.settings.VIDEO_TEMPLATE_PATH
    cfg_module.settings.VIDEO_TEMPLATE_PATH = str(TEMPLATE_PATH)
    try:
        pptx_path = generator._build_pptx(slides)
    finally:
        cfg_module.settings.VIDEO_TEMPLATE_PATH = original
    return Presentation(pptx_path)


def _placeholders(slide) -> dict[int, str]:
    """Return {idx: text} for all text-bearing placeholders in a slide."""
    result = {}
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        try:
            result[idx] = ph.text_frame.text
        except Exception:
            result[idx] = ""
    return result


def _notes(slide) -> str:
    return slide.notes_slide.notes_text_frame.text


# ------------------------------------------------------------------
# Slide count
# ------------------------------------------------------------------

def test_slide_count_matches_input(generator):
    slides = [
        TitleSlide(title="Intro", subtitle="Sub"),
        TitleAndContentSlide(title="Content", bullets=["A", "B"]),
        SectionHeaderSlide(title="Closing"),
    ]
    prs = build(generator, slides)
    assert len(prs.slides) == 3


def test_empty_slides_produces_empty_pptx(generator):
    prs = build(generator, [])
    assert len(prs.slides) == 0


# ------------------------------------------------------------------
# Layout selection
# ------------------------------------------------------------------

def test_title_slide_uses_correct_layout(generator):
    prs = build(generator, [TitleSlide(title="T", subtitle="S")])
    assert prs.slides[0].slide_layout.name == "Title Slide"


def test_title_and_content_uses_correct_layout(generator):
    prs = build(generator, [TitleAndContentSlide(title="T", bullets=["B"])])
    assert prs.slides[0].slide_layout.name == "Title and Content"


def test_section_header_uses_correct_layout(generator):
    prs = build(generator, [SectionHeaderSlide(title="T", text="X")])
    assert prs.slides[0].slide_layout.name == "Section Header"


def test_two_content_uses_correct_layout(generator):
    prs = build(generator, [TwoContentSlide(title="T", bullets_left=["L"], bullets_right=["R"])])
    assert prs.slides[0].slide_layout.name == "Two Content"


def test_comparison_uses_correct_layout(generator):
    prs = build(generator, [ComparisonSlide(title="T")])
    assert prs.slides[0].slide_layout.name == "Comparison"


def test_title_only_uses_correct_layout(generator):
    prs = build(generator, [TitleOnlySlide(title="T")])
    assert prs.slides[0].slide_layout.name == "Title Only"


def test_content_with_caption_uses_correct_layout(generator):
    prs = build(generator, [ContentWithCaptionSlide(title="T", bullets=["B"], caption="C")])
    assert prs.slides[0].slide_layout.name == "Content with Caption"


def test_picture_with_caption_uses_correct_layout(generator):
    prs = build(generator, [PictureWithCaptionSlide(title="T", image_prompt="img", caption="C")])
    assert prs.slides[0].slide_layout.name == "Picture with Caption"


# ------------------------------------------------------------------
# Placeholder content — TitleSlide
# ------------------------------------------------------------------

def test_title_slide_fills_title_and_subtitle(generator):
    prs = build(generator, [TitleSlide(title="Hello World", subtitle="A subtitle")])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Hello World"
    assert ph[1] == "A subtitle"


def test_title_slide_empty_subtitle(generator):
    prs = build(generator, [TitleSlide(title="Only Title", subtitle="")])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Only Title"


# ------------------------------------------------------------------
# Placeholder content — TitleAndContentSlide
# ------------------------------------------------------------------

def test_title_and_content_fills_title_and_bullets(generator):
    bullets = ["Point one", "Point two", "Point three"]
    prs = build(generator, [TitleAndContentSlide(title="My Title", bullets=bullets)])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "My Title"
    for bullet in bullets:
        assert bullet in ph[1]


def test_title_and_content_empty_bullets_skips_body(generator):
    prs = build(generator, [TitleAndContentSlide(title="T", bullets=[])])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "T"


# ------------------------------------------------------------------
# Placeholder content — SectionHeaderSlide
# ------------------------------------------------------------------

def test_section_header_fills_title_and_text(generator):
    prs = build(generator, [SectionHeaderSlide(title="Section A", text="Transition text")])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Section A"
    assert ph[1] == "Transition text"


# ------------------------------------------------------------------
# Placeholder content — TwoContentSlide
# ------------------------------------------------------------------

def test_two_content_fills_both_columns(generator):
    prs = build(generator, [TwoContentSlide(
        title="Compare",
        bullets_left=["Left 1", "Left 2"],
        bullets_right=["Right 1", "Right 2"],
    )])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Compare"
    assert "Left 1" in ph[1]
    assert "Right 1" in ph[2]


# ------------------------------------------------------------------
# Placeholder content — ComparisonSlide
# ------------------------------------------------------------------

def test_comparison_fills_headers_and_columns(generator):
    prs = build(generator, [ComparisonSlide(
        title="Comparison",
        header_left="Option A",
        bullets_left=["A1", "A2"],
        header_right="Option B",
        bullets_right=["B1", "B2"],
    )])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Comparison"
    assert ph[1] == "Option A"
    assert "A1" in ph[2]
    assert ph[3] == "Option B"
    assert "B1" in ph[4]


# ------------------------------------------------------------------
# Placeholder content — TitleOnlySlide
# ------------------------------------------------------------------

def test_title_only_fills_title(generator):
    prs = build(generator, [TitleOnlySlide(title="Big Statement")])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Big Statement"


# ------------------------------------------------------------------
# Placeholder content — ContentWithCaptionSlide
# ------------------------------------------------------------------

def test_content_with_caption_fills_all(generator):
    prs = build(generator, [ContentWithCaptionSlide(
        title="Title",
        bullets=["Bullet A", "Bullet B"],
        caption="Caption text",
    )])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Title"
    assert "Bullet A" in ph[1]
    assert ph[2] == "Caption text"


# ------------------------------------------------------------------
# Placeholder content — PictureWithCaptionSlide
# ------------------------------------------------------------------

def test_picture_with_caption_fills_title_and_caption(generator):
    prs = build(generator, [PictureWithCaptionSlide(
        title="Visual Slide",
        image_prompt="A futuristic city",
        caption="Fig 1. Futuristic city",
    )])
    ph = _placeholders(prs.slides[0])
    assert ph[0] == "Visual Slide"
    assert ph[2] == "Fig 1. Futuristic city"


def test_picture_with_caption_picture_placeholder_is_empty(generator):
    """Picture placeholder (idx=1) should be left empty — filled later by Gemini."""
    prs = build(generator, [PictureWithCaptionSlide(title="T", image_prompt="img", caption="C")])
    ph_indices = {p.placeholder_format.idx for p in prs.slides[0].placeholders}
    # idx=1 (picture) present in layout — text should be empty or not set
    if 1 in ph_indices:
        text = prs.slides[0].placeholders[1].text_frame.text if hasattr(
            prs.slides[0].placeholders[1], "text_frame"
        ) else ""
        assert text == ""


# ------------------------------------------------------------------
# Speaker notes
# ------------------------------------------------------------------

def test_speaker_notes_set(generator):
    prs = build(generator, [
        TitleSlide(title="T", subtitle="S", speaker_notes="This is the intro narration."),
    ])
    assert _notes(prs.slides[0]) == "This is the intro narration."


def test_no_speaker_notes_leaves_empty(generator):
    prs = build(generator, [TitleOnlySlide(title="T", speaker_notes="")])
    assert _notes(prs.slides[0]) == ""


# ------------------------------------------------------------------
# Mixed deck
# ------------------------------------------------------------------

def test_mixed_deck_layout_order(generator):
    slides = [
        TitleSlide(title="Intro", subtitle="Sub"),
        TitleAndContentSlide(title="Overview", bullets=["A", "B"]),
        TwoContentSlide(title="Side by Side", bullets_left=["L"], bullets_right=["R"]),
        ComparisonSlide(title="Versus", header_left="Pro", bullets_left=["P1"], header_right="Con", bullets_right=["C1"]),
        ContentWithCaptionSlide(title="Details", bullets=["D1"], caption="Note"),
        PictureWithCaptionSlide(title="Visual", image_prompt="sunset", caption="Sunset"),
        SectionHeaderSlide(title="Thank You", text="End"),
    ]
    prs = build(generator, slides)
    expected_layouts = [
        "Title Slide",
        "Title and Content",
        "Two Content",
        "Comparison",
        "Content with Caption",
        "Picture with Caption",
        "Section Header",
    ]
    assert len(prs.slides) == len(expected_layouts)
    for slide, expected in zip(prs.slides, expected_layouts):
        assert slide.slide_layout.name == expected


# ------------------------------------------------------------------
# Fallback: unknown layout name in template
# ------------------------------------------------------------------

def test_unknown_layout_skipped(generator, monkeypatch):
    """If the template doesn't have a layout, that slide is silently skipped."""
    import app.core.config as cfg_module
    monkeypatch.setattr(cfg_module.settings, "VIDEO_TEMPLATE_PATH", str(TEMPLATE_PATH))

    # Patch _SLIDE_TYPE_MAP to point TitleOnlySlide to a non-existent layout
    import app.services.video_generator as vg_module
    original = vg_module._SLIDE_TYPE_MAP["title_only"]
    vg_module._SLIDE_TYPE_MAP["title_only"] = (TitleOnlySlide, "Non Existent Layout")
    try:
        pptx_path = generator._build_pptx([
            TitleSlide(title="T", subtitle="S"),
            TitleOnlySlide(title="Skip me"),
            SectionHeaderSlide(title="End"),
        ])
    finally:
        vg_module._SLIDE_TYPE_MAP["title_only"] = original

    prs = Presentation(pptx_path)
    assert len(prs.slides) == 2  # TitleOnlySlide was skipped


# ------------------------------------------------------------------
# Fallback: missing template file
# ------------------------------------------------------------------

def test_missing_template_falls_back_to_blank(generator, monkeypatch):
    import app.core.config as cfg_module
    monkeypatch.setattr(cfg_module.settings, "VIDEO_TEMPLATE_PATH", "/nonexistent/template.pptx")

    pptx_path = generator._build_pptx([TitleSlide(title="T", subtitle="S")])
    prs = Presentation(pptx_path)
    # Should still produce 1 slide using blank presentation's layouts
    assert len(prs.slides) == 1
